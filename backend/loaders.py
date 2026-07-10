"""Document loaders — turn files of different formats into plain text.

Supported: .pdf, .docx, .pptx, .html/.htm, .txt/.md
Each loader takes a file path and returns one big string of text.
"""

from __future__ import annotations

from pathlib import Path


def load_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())


def load_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    # Drop non-content tags: scripts, styling, navigation menus.
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    return soup.get_text(separator="\n\n", strip=True)


def load_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes
                 if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if texts:
            slides.append("\n".join(texts))
    return "\n\n".join(slides)


def load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


# Maps a file extension to the loader that can read it.
LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".html": load_html,
    ".htm": load_html,
    ".txt": load_txt,
    ".md": load_txt,
}


def load_file(path: Path) -> str | None:
    """Load any supported file; returns None for unsupported types."""
    loader = LOADERS.get(path.suffix.lower())
    if loader is None:
        return None
    return loader(path)
